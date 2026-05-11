#!/usr/bin/env python3
"""MNN 推理流程 PPT — 匹配 ONNX Runtime PPT 风格 (深色背景)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
import os

DOCS = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(DOCS, "MNN_Inference_Flow_TechShare.pptx")

# ── ORT 同款配色 ──
BG       = RGBColor(0x1E, 0x1E, 0x1E)
BLUE     = RGBColor(0x00, 0x78, 0xD4)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x88, 0x88, 0x88)
MUTED    = RGBColor(0x60, 0x60, 0x60)
GREEN    = RGBColor(0x9C, 0xDC, 0xBE)
ORANGE   = RGBColor(0xE8, 0x6A, 0x17)
CARD_BG  = RGBColor(0x2D, 0x2D, 0x2D)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 工具 ──
def blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid(); bg_shape.fill.fore_color.rgb = BG
    bg_shape.line.fill.background()
    return s

def txt(slide, l, t, w, h, text, size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, name="Microsoft YaHei"):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = name; p.alignment = align
    return tf

def lines(slide, l, t, w, h, items, size=13, color=WHITE, gap=1.3):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        text, b = (item, False) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(size)
        p.font.bold = b; p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(size * (gap - 1))
    return tf

def rrect(slide, l, t, w, h, fill=CARD_BG, border=CARD_BG):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(0.5)
    return s

def code_block(slide, l, t, w, h, code_lines, size=11):
    s = rrect(slide, l, t, w, h)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(12); tf.margin_top = Pt(10)
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size); p.font.name = "Consolas"
        p.font.color.rgb = GREEN
        p.space_after = Pt(2)
    return s

def section_title(slide, num, title):
    txt(slide, 0.8, 0.3, 12, 0.7, f"{num}. {title}", 32, True, BLUE)

def footer(slide, text_str):
    txt(slide, 0.8, 7.1, 12, 0.4, text_str, 11, False, MUTED)

def arrow_right(slide, x, y, w=0.35, h=0.3):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = BLUE; s.line.fill.background()

def flow_card(slide, x, y, label, fill=BLUE, w=2.4):
    rrect(slide, x, y, w, 0.9, fill)
    txt(slide, x, y + 0.05, w, 0.8, label, 12, True, WHITE, PP_ALIGN.CENTER)


# ═══ Slide 1: 封面 ═══
s = blank()
txt(s, 1.5, 1.8, 10, 1.5, "MNN 推理引擎", 48, True, WHITE)
txt(s, 1.5, 3.3, 10, 1.0, "模型加载、图优化、后端调度 与 推理执行 全流程", 24, False, MUTED)
shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.0), Inches(4.2), Inches(5.3), Inches(0.02))
shape.fill.solid(); shape.fill.fore_color.rgb = BLUE; shape.line.fill.background()
txt(s, 1.5, 4.5, 10, 0.6, "基于 MNN 源码分析  |  FlatBuffers → Schedule → Pipeline → Backend → Execution", 14, False, GRAY)

# ═══ Slide 2: 目录 ═══
s = blank()
txt(s, 0.8, 0.4, 5, 0.7, "目  录", 36, True, BLUE)
toc = [
    ("1", "MNN 总体架构", "核心组件与层级"),
    ("2", "模型加载 — .mnn 文件 → Net + Op + Tensor", "FlatBuffers 零拷贝解析"),
    ("3", "Schedule — Pipeline 构建与 Op 映射", "createRuntime + 全部Op编入Pipeline + initConstTensors"),
    ("4", "图优化 — Pipeline::encode", "SizeComputer 形状推导 + GeometryComputer 算子分解"),
    ("5", "后端调度 — Pipeline::allocMemory", "Execution 创建 + 跨后端 Copy + useCount 内存复用"),
    ("6", "推理执行 — Pipeline::execute", "双层循环 onExecute + DIRECT/INDIRECT 模式"),
    ("7", "核心数据结构", "Command / Tensor / Backend / PipelineInfo"),
    ("8", "端到端全流程总结", "5 阶段流水线 + 核心设计模式"),
]
for i, (num, title, desc) in enumerate(toc):
    y = 1.2 + i * 0.58
    txt(s, 1.2, y, 0.6, 0.45, num, 22, True, BLUE)
    txt(s, 1.8, y, 6, 0.3, title, 15, True, WHITE)
    txt(s, 1.8, y + 0.28, 6, 0.25, desc, 11, False, MUTED)

# ═══ Slide 3: 总体架构 ═══
s = blank()
section_title(s, 1, "MNN 总体架构")
layers = [
    ("API 层", "Interpreter (createSession / resizeSession / runSession)  +  Module (Express 动态图)"),
    ("调度层", "Schedule 模型切分 → Pipeline 构建 → 全图 Tensor 映射 + 常量折叠"),
    ("管线层", "Pipeline: encode (图优化) → allocMemory (后端调度) → execute (推理执行)"),
    ("后端层", "Backend 抽象 (CPU / GPU / NPU) → Execution 多态 (onResize + onExecute)"),
]
for i, (name, desc) in enumerate(layers):
    y = 1.4 + i * 1.2
    rrect(s, 1.5, y, 2.2, 0.8, BLUE)
    txt(s, 1.5, y + 0.15, 2.2, 0.5, name, 16, True, WHITE, PP_ALIGN.CENTER)
    txt(s, 4.0, y + 0.1, 8.5, 0.6, desc, 13, False, WHITE)
footer(s, "核心设计：数学变换 (encode) 与硬件调度 (allocMemory) 两阶段分离 → 同一模型可在 CPU/GPU/NPU 多种后端运行，CPU 自动兜底")

# ═══ Slide 4: 模型加载 ═══
s = blank()
section_title(s, 2, "模型加载 — .mnn 文件 → Net + Op + Tensor")

# 横向流程
for i, (label, fill) in enumerate([
    ("model.mnn\n(FlatBuffers)", BLUE),
    ("FileLoader\n读文件头 + 验证魔数", CARD_BG),
    ("GetNet(buf)\n零拷贝反序列化", CARD_BG),
    ("Net + Op[] +\nTensor[] + Blob[]", CARD_BG),
]):
    x = 1.0 + i * 3.0
    flow_card(s, x, 1.5, label, fill)
    if i < 3:
        arrow_right(s, x + 2.55, 1.75)

# 左右代码块
code_block(s, 0.8, 3.0, 6.0, 3.8, [
    "// MNN.fbs — Net 定义 (schema/default/MNN.fbs:524)",
    "table Net {",
    "    bizCode: string;",
    "    extraTensorDescribe: [TensorDescribe];",
    "    extraInfo: ExtraInfo;",
    "    oplists: [Op];              // 算子列表",
    "    outputName: [string];       // 输出名列表",
    "    preferForwardType: ForwardType = CPU;",
    "    sourceType: NetSource = CAFFE;",
    "    tensorName: [string];       // 全图张量名",
    "    tensorNumber: int = 0;      // 张量总数",
    "    usage: Usage = INFERENCE;",
    "    subgraphs: [SubGraphProto];  // 子图",
    "    mnn_uuid: string;",
    "}",
    "",
    "table Op {",
    "    type: OpType;               // Conv / MatMul / ...",
    "    main: OpParameter;          // 算子参数",
    "    inputIndexes: [int];        // 输入 Tensor ID",
    "    outputIndexes: [int];       // 输出 Tensor ID",
    "    defaultDimCount: int;       // 维度数量",
    "}",
], 11)

code_block(s, 7.2, 3.0, 5.5, 3.8, [
    "// 关键技术点",
    "",
    "1. FlatBuffers 零拷贝",
    "   模型文件 = 一段连续内存",
    "   不需要 JSON/XML 解析",
    "   按偏移量直接访问字段",
    "   GetNet(buf) 一个指针操作",
    "",
    "2. NC4HW4 内存布局",
    "   Channel 维度按 4 对齐",
    "   为 ARM NEON 128-bit 设计",
    "   单指令同时处理 4×float32",
    "",
    "3. Resource 懒加载",
    "   权重不立即读入内存",
    "   首帧推理时按需加载",
    "   减少冷启动时间",
], 11)

footer(s, "FlatBuffers 是 MNN 区别于 ONNX(Protobuf) 的关键选择 → 解析极快，适合移动端启动性能要求高的场景")
# ═══ Slide 5: 使用实例 — Conv + ReLU ═══
s = blank()
section_title(s, 2, "模型加载实例 — Conv + ReLU")

# ── 左侧：模型逻辑结构 ──
rrect(s, 0.5, 1.15, 5.3, 0.45, BLUE)
txt(s, 0.7, 1.17, 4.9, 0.4, "模型逻辑结构", 15, True, WHITE, PP_ALIGN.CENTER)

# 2 个 Op 卡片 + Tensor 流转
for i, (label, fill) in enumerate([
    ("Conv2D\n3×3, S=1, P=1, 3ch→64ch\n持有 weight[64,3,3,3] + bias[64]", BLUE),
    ("ReLU\nf(x)=max(0,x)\n逐元素, 无参数", CARD_BG),
]):
    x = 0.8 + i * 2.6
    rrect(s, x, 1.85, 2.3, 1.05, fill)
    txt(s, x, 1.88, 2.3, 1.0, label, 9, False, WHITE, PP_ALIGN.CENTER)
    if i < 1:
        arrow_right(s, x + 2.35, 2.3, 0.2, 0.2)

# Tensor 流转表
rrect(s, 0.5, 3.1, 5.3, 3.5)
lines(s, 0.7, 3.15, 4.9, 3.4, [
    ("Tensor 数据流转 (inputIndexes / outputIndexes):", True),
    ("", False),
    ("Tensor[0]  input           NCHW [1, 3,   224, 224]   0.8 MB  (Usage: INPUT)", False),
    ("  │  Conv.inputIndexes  = [0, 2, 3]     ← input, weight, bias", False),
    ("  │  Conv.outputIndexes = [1]", False),
    ("Tensor[1]  conv_out        NCHW [1, 64,  112, 112]   0.8 MB  (Usage: NORMAL)", False),
    ("  │  ReLU.inputIndexes  = [1]", False),
    ("  │  ReLU.outputIndexes = [4]", False),
    ("Tensor[4]  relu_out        NCHW [1, 64,  112, 112]   0.8 MB  (Usage: OUTPUT)", False),
    ("", False),
    ("Tensor[2] conv_w   Blob float32s  [64,3,3,3] = 1728 floats  (Usage: CONSTANT)", False),
    ("Tensor[3] conv_b   Blob float32s  [64]                      (Usage: CONSTANT)", False),
    ("", False),
    ("2 Op, 4 Tensor (含权重), Conv 计算量 ~113 MFLOPS", False),
], 8.5, WHITE, 1.15)

# ── 右侧：FlatBuffers 加载后内存排布 ──
rrect(s, 6.2, 1.15, 6.6, 0.45, ORANGE)
txt(s, 6.4, 1.17, 6.2, 0.4, "FlatBuffers 加载后 — 连续 buffer 内存排布", 15, True, WHITE, PP_ALIGN.CENTER)

code_block(s, 6.2, 1.7, 6.6, 5.1, [
    '// FileLoader::load("model.mnn") → char* buf',
    '// GetNet(buf) → Net* (一个指针, 零拷贝)',
    '',
    'char* buf 连续内存:',
    '┌─────────────────────────────────────┐',
    '│ 0x1000  Net {                       │',
    '│   tensorName @ 0x2000               │',
    '│   oplists    @ 0x1A00               │',
    '│   outputName @ 0x2100               │',
    '│   tensorNumber = 5                  │',
    '│   usage = INFERENCE                 │',
    '├─────────────────────────────────────┤',
    '│ 0x1A00  Op[0] {                     │',
    '│   type = Conv2D, defaultDimCount=4  │',
    '│   inputIndexes:  [0, 2, 3]          │',
    '│   outputIndexes: [1]                │',
    '│   main @ 0x3000 → Convolution2D {   │',
    '│     kernelX=3, kernelY=3,           │',
    '│     strideX=1, strideY=1,           │',
    '│     padX=1, padY=1,                 │',
    '│     outputCount=64, group=1  }      │',
    '├─────────────────────────────────────┤',
    '│ 0x1A40  Op[1] {                     │',
    '│   type = ReLU                       │',
    '│   inputIndexes: [1], outputIndexes:[4] } │',
    '├─────────────────────────────────────┤',
    '│ 0x2000  tensorName:                 │',
    '│   "input" "conv_out" "conv_w" "bias"│',
    '│ 0x2100  outputName: [  "relu_out"  ]│',
    '├─────────────────────────────────────┤',
    '│ 0x3000  Convolution2D table         │',
    '│ 0x3200  Blob: float32s (conv 权重)  │',
    '│         64×3×3×3 = 1728 floats      │',
    '└─────────────────────────────────────┘',
], 6.5)

footer(s, "左: 逻辑有向图 (Op.inputIndexes/outputIndexes 描述数据流)   ↔   右: 物理连续 buffer (Net.tensorName/oplists → FlatBuffers offset 直接跳转, 零解析)")

# ═══ Slide 6: Schedule — 模型切分与 Pipeline 构建 ═══
s = blank()
section_title(s, 3, "Schedule — Pipeline 构建与 Op 映射")

# 4 步骤卡片横排
steps = [
    ("① createRuntime", "按 config.type\n创建 Backend\nCPU 始终创建 (兜底)\n→ RuntimeInfo", BLUE,
     "Schedule.cpp::schedule\nL110-180"),
    ("② 后端支持检查", "遍历每个 Op\nBackend::onCreate 试创建\n支持 → 加入当前 Pipeline\n不支持 → 切分新 Pipeline", ORANGE,
     "每个 config 独立\n生成一个 PipelineInfo"),
    ("③ initConstTensors", "遍历 CONSTANT Tensor\n必要时用 CPU 预计算\n结果写入 STATIC 内存\n后续 execute 直接读", CARD_BG,
     "Schedule::_scheduleUnit\nL290-330"),
    ("④ 构建 Session", "PipelineInfo[] → Session\n每个 Pipeline:\n  pair<BackendCache,\n        vector<OpCacheInfo>>\nSession 持有所有 Pipeline", GREEN,
     "Interpreter::createSession\nL180-220"),
]
for i, (title, desc, fill, note) in enumerate(steps):
    x = 0.3 + i * 3.25
    rrect(s, x, 1.3, 3.0, 0.55, fill)
    txt(s, x, 1.32, 3.0, 0.5, title, 13, True, WHITE, PP_ALIGN.CENTER)
    txt(s, x + 0.1, 2.0, 2.8, 2.2, desc, 10, False, WHITE)
    txt(s, x + 0.1, 4.3, 2.8, 0.8, note, 8, False, MUTED)

# 下半区: 源码逻辑
rrect(s, 0.5, 5.3, 12.3, 1.5)
lines(s, 0.7, 5.35, 11.9, 1.4, [
    ("Schedule::schedule 核心逻辑 (Schedule.cpp:291-331):", True),
    ("", False),
    ("  for (auto& config : configs) {", False),
    ("    auto oplists = _scheduleUnit(net, config, allTensors);   // \u2190 返回 \u2605\u6240\u6709\u2605 Op, 按原始顺序", False),
    ("    result.emplace_back({backendCache, oplists});            // \u2190 \u6bcf\u4e2a config \u4ea7\u751f\u4e00\u4e2a Pipeline, \u5305\u542b\u5168\u90e8 Op", False),
    ("  }", False),
    ("", False),
    ("\u5173\u952e: Schedule \u4e0d\u505a Op \u7c92\u5ea6\u5207\u5206! \u6bcf\u4e2a Pipeline \u5305\u542b\u5168\u90e8 Op, \u540e\u7aef\u5206\u914d\u5728 allocMemory._createExecutions \u9636\u6bb5\u9010 Op \u8bd5\u63a2", False),
], 9, WHITE, 1.15)

footer(s, "Schedule \u7684\u804c\u8d23: \u6309 config \u521b\u5efa Pipeline, \u6bcf\u4e2a Pipeline \u5305\u542b\u5168\u90e8 Op \u2192 \u540e\u7aef\u5206\u914d\u3001\u8de8\u540e\u7aef Copy \u7531 allocMemory \u9636\u6bb5\u5904\u7406")
# ═══ Slide 7: Schedule 实例 — 后端分配与执行顺序 ═══
s = blank()
section_title(s, 3, "Schedule 实例 — 同一 Pipeline, 逐 Op 后端分配")

# ── 上半区左: 逐 Op 后端试探 ──
rrect(s, 0.5, 1.2, 6.0, 0.4, BLUE)
txt(s, 0.7, 1.22, 5.6, 0.35, "allocMemory._createExecutions: 逐 Op 试探后端 (Pipeline.cpp)", 13, True, WHITE, PP_ALIGN.CENTER)

rrect(s, 0.5, 1.7, 6.0, 3.0)
lines(s, 0.7, 1.75, 5.6, 2.9, [
    ("原 Op 顺序: Conv \u2192 Softmax \u2192 BN \u2192 ReLU \u2192 Pool \u2192 FC", True),
    ("", False),
    ("_createExecutions 4 \u5c42\u4f18\u5148\u7ea7\u9010 Op \u8bd5\u63a2:", False),
    ("", False),
    ("Conv    \u2192 GPU.onCreate \u2713 \u2192 GPU Execution", False),
    ("Softmax \u2192 GPU.onCreate \u2717 \u2192 CPU \u515c\u5e95 \u2192 CPU Execution", False),
    ("BN      \u2192 GPU.onCreate \u2717 \u2192 CPU \u515c\u5e95 \u2192 CPU Execution", False),
    ("ReLU    \u2192 GPU.onCreate \u2713 \u2192 GPU Execution", False),
    ("Pool    \u2192 GPU.onCreate \u2713 \u2192 GPU Execution", False),
    ("FC      \u2192 GPU.onCreate \u2717 \u2192 CPU \u515c\u5e95 \u2192 CPU Execution", False),
    ("", False),
    ("\u7ed3\u679c: \u540c\u4e00\u4e2a Pipeline \u5185, GPU/CPU Execution \u6df7\u5408, \u987a\u5e8f\u4e0d\u53d8!", False),
], 10, WHITE, 1.2)

# ── 上半区右: Cross-backend Copy ──
rrect(s, 6.8, 1.2, 6.0, 0.4, ORANGE)
txt(s, 7.0, 1.22, 5.6, 0.35, "allocMemory._InsertCopy: 跨后端自动插入 Copy (Pipeline.cpp)", 13, True, WHITE, PP_ALIGN.CENTER)

rrect(s, 6.8, 1.7, 6.0, 3.0)
lines(s, 7.0, 1.75, 5.6, 2.9, [
    ("_InsertCopy \u68c0\u6d4b\u6bcf\u4e2a Execution \u7684 I/O Tensor \u5f52\u5c5e:", False),
    ("", False),
    ("Conv(GPU).output \u2192 Softmax(CPU).input", False),
    ("  \u2192 \u751f\u4ea7\u8005 GPU, \u6d88\u8d39\u8005 CPU, \u4e0d\u540c!", False),
    ("  \u2192 \u63d2\u5165 Copy Command (GPU\u2192CPU)", False),
    ("", False),
    ("BN(CPU).output \u2192 ReLU(GPU).input", False),
    ("  \u2192 \u63d2\u5165 Copy Command (CPU\u2192GPU)", False),
    ("", False),
    ("Pool(GPU).output \u2192 FC(CPU).input", False),
    ("  \u2192 \u63d2\u5165 Copy Command (GPU\u2192CPU)", False),
    ("", False),
    ("CONSTANT Tensor: \u7acb\u5373\u62f7\u8d1d; \u975e CONSTANT: \u63d2\u5165\u5ef6\u8fdf Copy", False),
    ("\u6bcf\u6b21 execute \u53ea\u62f7\u8d1d\u6709\u810f\u6807\u8bb0\u7684\u6570\u636e", False),
], 9, WHITE, 1.2)

# ── 下半区: 执行顺序 + Session ──
rrect(s, 0.5, 5.0, 12.3, 1.8)
lines(s, 0.7, 5.05, 11.9, 1.7, [
    ("Pipeline::execute \u6267\u884c\u987a\u5e8f (\u540c\u4e00\u4e2a Pipeline, \u4e00\u8d9f\u4e32\u884c):", True),
    ("", False),
    ("  Conv(GPU) \u2192 [Copy GPU\u2192CPU] \u2192 Softmax(CPU) \u2192 BN(CPU) \u2192 [Copy CPU\u2192GPU] \u2192 ReLU(GPU) \u2192 Pool(GPU) \u2192 [Copy GPU\u2192CPU] \u2192 FC(CPU)", False),
    ("", False),
    ("\u987a\u5e8f\u4fdd\u8bc1\u4e09\u91cd\u673a\u5236:", False),
    ("  1. \u539f\u59cb Op \u6570\u7ec4\u987a\u5e8f (generateScheduleGraph \u6309 net->oplists() \u904d\u5386, \u4e0d\u6253\u4e71)", False),
    ("  2. initPipelineInfosFromOps \u6309 ops \u987a\u5e8f\u8ffd\u52a0 OpCacheInfo, \u4e0d\u4e22\u5f03\u4efb\u4f55 Op", False),
    ("  3. Pipeline::execute \u5185\u5c42 for \u5faa\u73af\u4e32\u884c, Tensor \u6570\u636e\u4f9d\u8d56\u81ea\u7136\u4e32\u8d77\u6267\u884c\u987a\u5e8f", False),
], 9, WHITE, 1.15)

footer(s, "\u6b63\u786e\u7406\u89e3: \u4e0d\u662f\u5207\u5206\u591a\u4e2a Pipeline, \u800c\u662f\u540c\u4e00 Pipeline \u5185\u9010 Op \u540e\u7aef\u5206\u914d + \u81ea\u52a8 Copy \u4e2d\u8f6c \u2192 \u539f\u59cb\u987a\u5e8f\u5929\u7136\u4fdd\u8bc1")



# ═══ Slide 8: 图优化 ═══
s = blank()
section_title(s, 4, "图优化 — Pipeline::encode (纯数学变换, 不依赖硬件)")

for i, (title, desc, color) in enumerate([
    ("形状推导", "70+ SizeComputer 子类\n每个算子一个数学公式\nConv: Oh=(Ih-Kh+2P)/S+1\nPool/MatMul/BinaryOp ...\n找不到 → 拷贝 inputs[0] 形状", BLUE),
    ("算子分解", "高阶 Op → 低阶 Command\nConv → Im2Col + MatMul\nConcat → Region → Raster\n逐元素 Op → 1:1 透传\n(等待 opFuse 融合)", ORANGE),
    ("常量折叠 & 融合", "CONSTANT Op 提前 CPU 执行\n结果写入 STATIC 内存\n脏标记链避免重复计算\nopFuse: GPU 逐元素融合\n(MNN_BUILD_CODEGEN)", GREEN),
]):
    x = 0.8 + i * 4.2
    rrect(s, x, 1.4, 3.8, 0.6, color)
    txt(s, x, 1.42, 3.8, 0.55, title, 16, True, WHITE, PP_ALIGN.CENTER)
    lines(s, x + 0.2, 2.2, 3.4, 4.2, desc.split("\n"), 12, WHITE, 1.6)

rrect(s, 0.8, 6.3, 11.8, 0.6, BLUE)
txt(s, 1.0, 6.35, 11.4, 0.5,
    "OpResizeCache 增量检测: match(inputs) 逐项比对上次输入形状 → 形状未变则跳过 SizeCompute → 大幅减少动态模型 resize 开销",
    11, False, WHITE)

footer(s, "encode 是 MNN 独有的「图编译」阶段 → ONNX Runtime 的 Graph::Resolve + GraphOptimizer 的混合体")
# ═══ Slide 9: 实例 — Conv+BN+ReLU 的 encode 过程 ═══
s = blank()
section_title(s, 4, "图优化实例 — Conv+BN+ReLU encode")

# ── 上半区：Before → After 横排 ──
rrect(s, 0.5, 1.2, 5.0, 0.45, BLUE)
txt(s, 0.7, 1.22, 4.6, 0.4, "encode 前: 原始模型图", 14, True, WHITE, PP_ALIGN.CENTER)

# Before: 3 个 Op 卡片
for i, (label, fill) in enumerate([
    ("Conv2D\nk=3,s=2,p=1\n3ch→64ch", BLUE),
    ("BatchNorm\nmean,var\nscale,bias", CARD_BG),
    ("ReLU", CARD_BG),
]):
    x = 0.6 + i * 1.6
    rrect(s, x, 1.85, 1.4, 0.95, fill)
    txt(s, x, 1.88, 1.4, 0.9, label, 8, False, WHITE, PP_ALIGN.CENTER)
    if i < 2:
        arrow_right(s, x + 1.43, 2.22, 0.14, 0.16)

# 中间大箭头
txt(s, 5.6, 2.1, 1.2, 0.5, "encode\n───>", 14, True, ORANGE, PP_ALIGN.CENTER)

# After: encode 后三个阶段的卡片
rrect(s, 7.0, 1.2, 5.8, 0.45, ORANGE)
txt(s, 7.2, 1.22, 5.4, 0.4, "encode 后: 3 步变换结果", 14, True, WHITE, PP_ALIGN.CENTER)

after_items = [
    ("1. 形状推导", "ConvSizeComputer:\nOh=(224-3+2)/2+1=112\n→ [1,64,112,112]\nBN/ReLU: 逐层推导\n所有 Tensor 形状确定", "SizeComputerSuite\n::search(OpType)"),
    ("2. 算子分解", "Conv2D → 3 个 Command:\n① Im2Col (数据重排, 0 FLOP)\n② MatMul (21.7M MAC)\n③ BiasAdd (加偏置)\nFLOPs 不变, GEMM 密集计算\nNEON 单指令 4×float", "GeometryComputer\n::compute"),
    ("3. 常量折叠 & 融合", "BN 是 CONSTANT Op:\n  CPU 提前算 mean/var/scale\n  折叠进 Conv 权重/偏置\n  → BN 执行时直接跳过!\nReLU → Conv.relu=true 融合", "encode 阶段的\nCPU 预计算"),
]
for i, (title, desc, note) in enumerate(after_items):
    x = 7.1 + i * 1.95
    rrect(s, x, 1.85, 1.8, 2.3)
    txt(s, x + 0.05, 1.88, 1.7, 0.3, title, 9, True, WHITE)
    txt(s, x + 0.05, 2.2, 1.7, 1.6, desc, 7, False, WHITE)
    txt(s, x + 0.05, 3.8, 1.7, 0.3, note, 6.5, False, MUTED)

# ── 下半区：结果对比 ──
rrect(s, 0.5, 4.4, 12.3, 2.4)
lines(s, 0.7, 4.45, 11.9, 2.3, [
    ("encode 的收益:", True),
    ("", False),
    ("原始 (逐 Op 执行):", False),
    ("  [Conv2D] → 中间结果写回内存 → [BatchNorm] → 读+算+写回 → [ReLU] → 读+算+写回", False),
    ("  每个 Op 各自读输入、写输出, 特征图 (64\u00d7112\u00d7112=3.2MB) 被反复搬运", False),
    ("", False),
    ("encode 后:", False),
    ("  算子分解: 直接卷积 21.7M MAC → Im2Col+GEMM 同样 21.7M MAC, 但 GEMM 密集 + NEON SIMD → ~4× 实际执行提速", False),
    ("  常量折叠: BN 参数预先融进 Conv 权重, 消除 BN 的读写 → 省 2 次特征图遍历 (~6.4MB 内存访问)", False),
    ("  算子融合: ReLU 内联到 Conv (relu=true), 消除 ReLU 的读写 → 再省 2 次特征图遍历 (~6.4MB 内存访问)", False),
    ("", False),
    ("合计: FLOPs 不变 (43.4M), 但 GEMM 密集 + SIMD → ~4× 加速; 折叠+融合 消除 ~12.8MB 内存访问; 整体推理延迟大幅下降!", False),
], 9, WHITE, 1.15)

footer(s, "encode 的威力: 数学变换完全独立于硬件, 编译期做完所有「纸上推演」→ allocMemory 只需关心内存和执行, execute 只需跑最简 Command")


# ═══ Slide 10: 后端调度 ═══
s = blank()
section_title(s, 5, "后端调度 — Pipeline::allocMemory (硬件相关)")

for i, (title, subtitle, desc, note) in enumerate([
    ("阶段1: _createExecutions", "4层优先级",
     "1) executionCache (O(1)命中)\n2) KV Cache 克隆\n3) 主 Backend::onCreate\n4) CPU 兜底",
     "Raster / BinaryOp\n不缓存 (形状敏感)"),
    ("阶段2: _SetTensorBackend", "标记归属",
     "遍历 Execution\n标记 I/O Tensor 的\n所属 Backend\n跳过 Copy 算子的 Tensor",
     ""),
    ("阶段3: _InsertCopy", "跨后端 Wrap",
     "needWrap() 检测跨后端\nCONSTANT → 立即拷贝\n非常量 → 插入 Copy Command\nworkInputs 在此被替换",
     "每次 execute 只拷贝\n有脏标记的数据"),
    ("阶段4: _allocForTensor", "useCount 内存复用",
     "统计引用计数\n归零 → 申请内存\n递减 → 释放复用\n三种类型: STATIC /\nDYNAMIC_SEPERATE / DYNAMIC",
     "峰值内存 = 最大同时\n存活 Tensor 之和"),
]):
    x = 0.3 + i * 3.25
    rrect(s, x, 1.4, 3.0, 0.6, BLUE)
    txt(s, x, 1.42, 3.0, 0.55, title, 12, True, WHITE, PP_ALIGN.CENTER)
    txt(s, x + 0.1, 2.1, 2.8, 0.3, subtitle, 14, True, ORANGE)
    lines(s, x + 0.1, 2.5, 2.8, 3.5, desc.split("\n"), 11, WHITE, 1.5)
    if note:
        lines(s, x + 0.1, 5.6, 2.8, 1.0, note.split("\n"), 10, MUTED)

footer(s, "allocMemory 把 encode 的数学结果「落地」到具体硬件和内存 → SessionState::FinalizeSessionState 的对应阶段")

# ═══ Slide 11: 推理执行 ═══
s = blank()
section_title(s, 6, "推理执行 — Pipeline::execute")

code_block(s, 0.5, 1.4, 6.5, 5.0, [
    "// Pipeline::execute() — Session::run → for each Pipeline",
    "ErrorCode Pipeline::execute() {",
    "    _copyInputs();          // ① host→device 输入拷贝",
    "    _enterExecute();        // ② Backend::onExecuteBegin",
    "",
    "    for (auto& info : mInfo.second) {     // 遍历 OpCacheInfo",
    "        if (CONSTANT == info.type) continue; // 跳过 (已折叠)",
    "",
    "        for (auto& cmd : info.executeBuffer.command) {",
    "            // ★ 虚函数多态 — 核心执行点",
    "            cmd->execution->onExecute(",
    "                cmd->workInputs,   // 物理输入 (_InsertCopy 后)",
    "                cmd->workOutputs); // 物理输出",
    "        }",
    "    }",
    "    _exitExecute();          // ③ Backend::onExecuteEnd",
    "}",
], 11)

rrect(s, 7.5, 1.4, 5.3, 2.3)
lines(s, 7.7, 1.5, 5.0, 2.1, [
    ("两种执行模式:", True),
    ("", False),
    ("DIRECT (CPU 默认):", True),
    ("  onExecute 被调用 → 直接计算完成", False),
    ("  例: CPU 卷积, 进出函数即算完", False),
    ("", False),
    ("INDIRECT (GPU 使用):", True),
    ("  onExecute → 录制命令到队列", False),
    ("  onExecuteBegin → 批量提交 GPU", False),
    ("  onExecuteEnd → 等待完成/取回结果", False),
    ("  好处: 减少 CPU-GPU 通信开销", False),
], 11, WHITE)

rrect(s, 7.5, 3.9, 5.3, 2.5)
lines(s, 7.7, 4.0, 5.0, 2.3, [
    ("Command 结构:", True),
    ("", False),
    ("struct Command {", False),
    ("  const Op* op;            // 算子描述", False),
    ("  vector<Tensor*> inputs;  // 逻辑引用 (不变)", False),
    ("  vector<Tensor*> outputs;", False),
    ("  vector<Tensor*> workInputs; // 物理执行", False),
    ("  vector<Tensor*> workOutputs;", False),
    ("  shared_ptr<Execution> execution;", False),
    ("  int group;               // 内存分组标记", False),
    ("};", False),
], 11, GREEN)

footer(s, "execute 是「真正跑计算」的地方 — Session 持有多个 Pipeline, 每个 Pipeline 持有 N 个 Command, 每个 Command 持有 1 个 Execution")

# ═══ Slide 12: 核心数据结构 ═══
s = blank()
section_title(s, 7, "核心数据结构")

for j, (text, x, w) in enumerate([
    ("数据结构", 0.5, 2.5), ("定义位置", 3.0, 4.0), ("职责", 7.0, 5.8),
]):
    rrect(s, x, 1.4, w, 0.5, BLUE)
    txt(s, x, 1.42, w, 0.45, text, 13, True, WHITE, PP_ALIGN.CENTER)

rows = [
    ("Net / Op / Tensor", "schema/default/MNN.fbs", "模型格式: 算子类型 + 参数 + I/O ID + 权重"),
    ("Interpreter", "include/MNN/Interpreter.hpp", "门面 API: createSession / resizeSession / runSession"),
    ("Schedule", "source/core/Schedule.cpp", "模型切分: createRuntime → initTensors → _scheduleUnit"),
    ("Session", "source/core/Session.cpp", "会话管理: mPipelines + mNeedResize/mNeedMalloc 两阶段标记"),
    ("Pipeline", "source/core/Pipeline.cpp", "管线引擎: encode + allocMemory + execute 三阶段"),
    ("Backend", "source/core/Backend.hpp", "后端抽象: onCreate (工厂) + onAcquireBuffer (内存)"),
    ("Execution", "source/core/Execution.hpp", "算子执行器: virtual onResize() + onExecute() = 0"),
    ("SizeComputer", "source/shape/SizeComputer.cpp", "形状推导: 70+ 子类 × 每个算子的数学公式"),
    ("Command", "source/core/Command.hpp", "执行单元: op + inputs/outputs + workInputs + execution"),
    ("OpCacheInfo", "source/core/Schedule.hpp", "运行时缓存: executeBuffer + computeCache + executionCache"),
    ("PipelineInfo", "source/core/Schedule.hpp", "子图结构: pair<BackendCache, vector<OpCacheInfo>>"),
    ("RuntimeInfo", "source/core/Session.hpp", "多后端环境: map<Type, Runtime> + 默认 CPU Runtime"),
]

for i, (name, location, desc) in enumerate(rows):
    y = 2.0 + i * 0.42
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x35, 0x35, 0x35)
    for (text_val, l, w) in [(name, 0.5, 2.5), (location, 3.0, 4.0), (desc, 7.0, 5.8)]:
        rrect(s, l, y, w, 0.39, bg)
        txt(s, l + 0.1, y + 0.05, w - 0.2, 0.30, text_val, 10, False, WHITE if text_val == name else MUTED)

footer(s, "MNN 的核心数据结构围绕「一次推理的执行单元」展开 → Command 是原子操作，Pipeline 是执行容器，Session 是生命周期管理者")

# ═══ Slide 13: 端到端全流程总结 ═══
s = blank()
section_title(s, 8, "端到端全流程总结")

# 横向流程
for i, (label, fill) in enumerate([
    ("model.mnn", BLUE), ("FlatBuffers\n解析", CARD_BG), ("Schedule\n切分子图", CARD_BG),
    ("encode\n图优化", CARD_BG), ("allocMemory\n后端调度", CARD_BG), ("execute\n推理执行", CARD_BG),
    ("结果\nTensor", CARD_BG),
]):
    x = 0.3 + i * 1.85
    flow_card(s, x, 1.5, label, fill, w=1.55)
    if i < 6:
        arrow_right(s, x + 1.65, 1.68, 0.18, 0.25)

# 4列设计模式
for i, (title, desc, color) in enumerate([
    ("注册-查找 分离", "SizeComputerSuite\nsearch(OpType) O(1)\nBackend::onCreate\n工厂模式多态\nExecutionCache 复用", BLUE),
    ("图与执行解耦", "encode 只看数学\n不碰硬件 Backend\nallocMemory 只看硬件\n不碰形状逻辑\n静态模型跳过 Geometry", ORANGE),
    ("内存生命周期管理", "useCount 引用计数\n每次 allocMemory 重算\n归零即释放\n三种类型分层\nfixResizeCache 保护", GREEN),
    ("多层 Fallback", "CPU 始终兜底\nGPU 不支持→自动回退\n异步调优 + 缓存\n_InsertCopy 自动\n处理跨后端传输", BLUE),
]):
    x = 0.3 + i * 3.25
    rrect(s, x, 2.7, 3.0, 0.6, color)
    txt(s, x, 2.72, 3.0, 0.55, title, 13, True, WHITE, PP_ALIGN.CENTER)
    lines(s, x + 0.1, 3.5, 2.8, 3.0, desc.split("\n"), 11, WHITE)

# 底部注册/查找示意
code_block(s, 0.3, 6.5, 6.3, 0.8, [
    "// 注册 (编译期)                         // 查找 (运行时)",
    "REGISTER_SHAPE(OpType_Conv, ConvSizeComputer); → SizeComputerSuite::search(OpType_Conv)",
    "Backend::onCreate(inputs, outputs, op);        → ExecutionCache::find(op)",
], 10)

code_block(s, 6.9, 6.5, 5.8, 0.8, [
    "// MNN 独有的 encode/allocMemory 两阶段设计",
    "encode:     纯数学, 不依赖硬件, 可离线静态化",
    "allocMemory: 纯硬件, 不关心形状, 专注于性能",
], 10)

# ═══ Slide 14: Thank You ═══
s = blank()
txt(s, 1.5, 2.0, 10, 1.5, "Thank You", 48, True, WHITE, PP_ALIGN.CENTER)
txt(s, 1.5, 3.5, 10, 1.0, "Q & A", 24, False, MUTED, PP_ALIGN.CENTER)
shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.2), Inches(4.8), Inches(5.0), Inches(0.02))
shape.fill.solid(); shape.fill.fore_color.rgb = BLUE; shape.line.fill.background()
txt(s, 1.5, 5.3, 10, 0.6, "github.com/alibaba/MNN", 14, False, GRAY, PP_ALIGN.CENTER)
txt(s, 1.5, 5.8, 10, 0.6, "基于 MNN 源码分析 | FlatBuffers → Schedule → Pipeline → Backend → Execution", 12, False, GRAY, PP_ALIGN.CENTER)

# ── 保存 ──
prs.save(OUT)
print(f"PPT saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
